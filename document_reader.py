"""
Extracts raw text content from PDF, Word (.docx), Excel (.xlsx), and PowerPoint (.pptx) files.

This is deliberately dumb/mechanical - it just gets everything readable
out of the file as plain text. The actual understanding (what's the tour
name, what's included, what are the destinations, etc.) happens in the
next step (ai_extractor.py), which is much better suited to messy,
inconsistent supplier documents than hand-written parsing rules.

TABLE STRUCTURE IS NOT DECORATION. A rate sheet carries its meaning in the GRID: which price
sits under which season, which header spans which columns. Flattening a row to "a | b | c"
throws that away, and no amount of prompting recovers it afterwards.

CONFIRMED REAL FAILURE (product owner, Le Fayan Fleet 2026/27): "the document reader is doing
very difficult to read those styles of documents... after 10 times trying and even with TELL AI
I was unable to publish that." That rate table has a season header spanning two columns, three
stacked date ranges under one season, From/To sub-headers on different rows for different
seasons, and - worst of all - Word stores the prices as a SEPARATE table with no header row at
all. The old renderer emitted a merged cell once per column it spanned, so "Season 2026 / 2027"
appeared eight times and every price appeared twice, and the header table and the price table
arrived as two unrelated blocks. The model had no way to know which $565 belonged to which
season - and no way to say so either, so it simply produced something wrong.

This module now preserves three things a model cannot reconstruct on its own:
  - WHICH COLUMNS a merged cell spans, written explicitly instead of by repetition
  - a COLUMN RULER, so any cell can be referred to by position
  - whether a table is a CONTINUATION of the one above it, sharing its headers
"""
import os

# Stamped on every delivery. app.py compares this against its own build string and says
# so on screen when they differ - a partial push (one file committed, another not) used to
# surface only as a traceback whose line numbers pointed at unrelated code.
MODULE_BUILD = "2026-09-04-pptx-text-and-image-extraction"

_EMPTY_CELL = "·"          # visible placeholder, so a blank column is not silently swallowed


def _render_grid(rows, table_label, previous_column_count=None, previous_paths=None,
                 previous_label=None):
    """Render one table as text that keeps its grid.

    `rows` is a list of rows, each a list of (text, start_col, end_col) spans - one entry per
    DISTINCT cell, carrying the columns it covers, rather than one entry per column.

    Returns (lines, column_count, column_paths)."""
    if not rows:
        return [], 0, []

    width = max((span[2] + 1) for row in rows for span in row) if any(rows) else 0
    lines = [f"[{table_label}]",
             "COLUMNS: " + " | ".join(f"C{i + 1}" for i in range(width))]

    # A table whose first row is entirely empty and which matches the previous table's width is
    # almost certainly a continuation: in Word, one visual table is often stored as two, and the
    # second half then has no headers at all.
    first_row_blank = all(not (span[0] or "").strip() for span in rows[0])
    is_continuation = bool(first_row_blank and previous_column_count == width and width > 1)
    if is_continuation:
        lines.append(
            f"NOTE: this table has NO header row of its own, and has the same number of columns "
            f"as the table immediately above it. Its columns line up with that table's headers "
            f"one for one - read C1..C{width} there to know what each column here means. The "
            f"BY COLUMN list below already has those headings merged in.")

    for row_index, row in enumerate(rows, start=1):
        cells = []
        for text, start, end in row:
            text = (text or "").strip() or _EMPTY_CELL
            if end > start:
                cells.append(f"{text} «spans C{start + 1}-C{end + 1}»")
            else:
                cells.append(f"{text} «C{start + 1}»")
        lines.append(f"R{row_index}: " + " | ".join(cells))

    # The column-wise view. Capped, because on a very wide table it would double the text for
    # no gain - and a 40-column rate sheet does not exist.
    paths = _column_paths(rows, width)
    if 1 < width <= _MAX_COLUMN_VIEW_COLUMNS:
        lines.append("BY COLUMN (the same table read downwards - use this to match a value to "
                     "its heading, rather than counting across):")
        for index, path in enumerate(paths):
            inherited = []
            if is_continuation and previous_paths and index < len(previous_paths):
                inherited = previous_paths[index]
            full = inherited + [v for v in path if v]
            if not full:
                continue
            lines.append(f"  C{index + 1} = " + " > ".join(full))
    lines.append(f"[/{table_label}]")
    return lines, width, paths


_MAX_COLUMN_VIEW_COLUMNS = 40


def _cell_at(row, column):
    for text, start, end in row:
        if start <= column <= end:
            return (text or "").strip()
    return ""


def _column_paths(rows, width):
    """The same table read DOWNWARDS - one labelled path per column.

    This is the "translate it into a structured list first" step, done deterministically here
    instead of asked of the model. Reading a season grid ACROSS means holding merged spans,
    sub-headers on different rows and stacked date ranges in mind at once. Reading it DOWN means
    nothing harder than following a list, so C2 becomes:

        Season 2026 / 2027 > Normal > From > 24/9/2026 > 7/1/2027 > 5/4/2027

    and on the price table below it, with each figure carrying the row it came from:

        ... > Single Luxury Cabin: $565 > Per person in Double Luxury Cabin: $353 > ...

    The association between a season, its dates and the price of a particular cabin stops being
    an inference the model has to make and becomes something it can simply read."""
    if not rows or width < 1:
        return []

    # Column 0 is a row label only when it actually varies down the table. On a header block it
    # is one word merged down the side ("Seasonality"), and prefixing every value with it would
    # add noise rather than meaning.
    first_column = [_cell_at(row, 0) for row in rows]
    filled = [v for v in first_column if v]
    use_row_labels = len(set(filled)) > 1 and len(set(filled)) >= max(2, len(filled) // 2)

    paths = []
    for column in range(width):
        values = []
        for row_index, row in enumerate(rows):
            text = _cell_at(row, column)
            if not text:
                continue
            # Skip a repeat of the value directly above: a label merged vertically down the side
            # adds nothing by being said six times.
            if values and values[-1].endswith(text) and not use_row_labels:
                continue
            label = first_column[row_index] if (use_row_labels and column > 0) else ""
            if label and label != text:
                values.append(f"{label}: {text}")
            elif not values or values[-1] != text:
                values.append(text)
        paths.append(values)
    return paths


def _docx_row_spans(row):
    """Distinct cells in a Word table row, with the columns each one covers.

    python-docx returns the SAME underlying <w:tc> element once per column a merged cell covers,
    so identity is what reveals the merge. Comparing text instead would wrongly merge two
    genuinely separate columns that happen to hold the same number - and on a rate sheet, two
    seasons at the same price is completely ordinary."""
    spans = []
    previous_tc = None
    for index, cell in enumerate(row.cells):
        tc = cell._tc
        if tc is previous_tc and spans:
            spans[-1][2] = index
        else:
            spans.append([cell.text.strip(), index, index])
            previous_tc = tc
    return [tuple(s) for s in spans]


def extract_text_from_pdf(file_path: str) -> str:
    import pdfplumber

    parts = []
    previous_pdf_width = None
    previous_pdf_paths = None
    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            parts.append(f"--- Page {page_num} ---\n{text}")

            for table_index, table in enumerate(page.extract_tables(), start=1):
                # pdfplumber gives a merged cell's text once and None for the columns it covers,
                # so a run of Nones after a value IS the span. Reconstructing it here means a PDF
                # rate sheet reads the same way a Word one does.
                rows = []
                for raw_row in table:
                    spans, last = [], None
                    for index, cell in enumerate(raw_row):
                        text = (cell or "").strip()
                        if text:
                            spans.append([text, index, index])
                            last = spans[-1]
                        elif last is not None and cell is None:
                            last[2] = index
                        else:
                            spans.append(["", index, index])
                            last = None
                    rows.append([tuple(s) for s in spans])
                lines, width, paths = _render_grid(rows, f"TABLE p{page_num}.{table_index}",
                                                    previous_pdf_width, previous_pdf_paths)
                parts.extend(lines)
                previous_pdf_width, previous_pdf_paths = width, paths

    return "\n".join(parts)


def extract_text_from_docx(file_path: str) -> str:
    """Paragraphs and tables IN DOCUMENT ORDER.

    CONFIRMED REAL DEFECT (Le Fayan Fleet): python-docx exposes doc.paragraphs and doc.tables as
    two separate collections, and reading them one after the other put every table at the END of
    the text. In that document the rate table sits directly under "Below are the net agent
    rates:" but arrived after thirty-one paragraphs, the last of which price Abu Simbel and a
    balloon ride. A model reading that sees a table of numbers whose nearest context is a
    completely different price list - so the sentence that says these rates are "per suite per
    night" is nowhere near the rates it describes.
    """
    import docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = docx.Document(file_path)
    parts = []
    previous_width = None
    previous_paths = None
    table_num = 0

    for child in doc.element.body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            text = Paragraph(child, doc).text.strip()
            if text:
                parts.append(text)
        elif tag == "tbl":
            table_num += 1
            table = Table(child, doc)
            try:
                rows = [_docx_row_spans(row) for row in table.rows]
            except Exception:
                # A table python-docx cannot walk (odd nesting, corrupt merge) must not cost the
                # whole document - fall back to the flat rendering for that one table.
                parts.append(f"[TABLE {table_num}]")
                for row in table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
                parts.append(f"[/TABLE {table_num}]")
                previous_width, previous_paths = None, None
                continue
            lines, width, paths = _render_grid(rows, f"TABLE {table_num}", previous_width,
                                                previous_paths)
            parts.extend(lines)
            previous_width, previous_paths = width, paths

    return "\n".join(parts)


def extract_text_from_pptx(file_path: str) -> str:
    """Slide text (title/body/text-box shapes, in shape order, including inside grouped shapes)
    plus speaker notes, with a simple flattened rendering for any table shapes.

    CONFIRMED REAL GAP (Chris, 2026-09-04): a supplier can send a whole tour/excursion writeup as
    a PowerPoint deck (e.g. "1 Day Diving Course.pptx") - this module raised "Unsupported file
    type" for '.pptx' before this fix, so the document couldn't even be uploaded. PowerPoint
    decks from suppliers are typically photo-heavy marketing material, not rate-grid tables like
    the Word/Excel documents this module was originally built to survive - so unlike
    extract_text_from_docx/extract_text_from_xlsx, this does NOT do merge-aware grid
    reconstruction (_render_grid) for tables yet. If a real pptx rate table with merged cells
    turns up, that's the trigger to bring the same treatment here."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    parts = []

    def _walk_shapes(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _walk_shapes(shape.shapes)
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts.append(f"[SLIDE {slide_num}]")
        _walk_shapes(slide.shapes)
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                parts.append(f"[SLIDE {slide_num} NOTES] {notes}")

    return "\n".join(parts)


def extract_text_from_xlsx(file_path: str) -> str:
    """Sheets, grid-preserving - the same span/ruler treatment extract_text_from_pdf and
    extract_text_from_docx already give their tables.

    CONFIRMED BUG FIX (audit 2026-09-01, MEDIUM/LOW batch 2): this used to just join each row's
    cell values with " | ", with no awareness of merged cells at all - the exact left-to-right
    guessing this module's own module docstring says the grid-preservation mechanism exists to
    prevent. A merged season header spanning several columns came out as one value floating
    above a row of blanks, with nothing telling the model which price columns it actually
    covers - forcing it into exactly the left-to-right column counting the prompts explicitly
    forbid. Excel rate sheets are common enough (Excel is the natural format for a tariff grid)
    that this was a real, not theoretical, gap."""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]

        # Map every merged region so a season header spanning several columns becomes ONE cell
        # covering those columns (matching _docx_row_spans/the PDF reader's span reconstruction),
        # instead of being flattened into a repeated-or-blank run of individual cells.
        merge_top_left = {}    # (row, col), 1-indexed -> (min_col, max_col) of its merge
        merge_covered = set()  # every OTHER (row, col) inside a merged range, to skip entirely
        for merged_range in sheet.merged_cells.ranges:
            merge_top_left[(merged_range.min_row, merged_range.min_col)] = (
                merged_range.min_col, merged_range.max_col)
            for r in range(merged_range.min_row, merged_range.max_row + 1):
                for c in range(merged_range.min_col, merged_range.max_col + 1):
                    if (r, c) != (merged_range.min_row, merged_range.min_col):
                        merge_covered.add((r, c))

        rows = []
        for row_cells in sheet.iter_rows():
            spans = []
            row_has_content = False
            for cell in row_cells:
                coord = (cell.row, cell.column)
                if coord in merge_covered:
                    continue
                text = "" if cell.value is None else str(cell.value)
                text = text.strip()
                if text:
                    row_has_content = True
                start = end = cell.column - 1
                if coord in merge_top_left:
                    end = merge_top_left[coord][1] - 1
                spans.append((text, start, end))
            if row_has_content:
                rows.append(spans)

        if not rows:
            continue
        parts.append(f"--- Sheet: {sheet_name} ---")
        # Each sheet is its own independent table, not a continuation of the previous one, so no
        # previous_column_count/previous_paths are carried across sheets.
        lines, _, _ = _render_grid(rows, f"SHEET {sheet_name}")
        parts.extend(lines)

    return "\n".join(parts)


_MIN_USEFUL_CHARS = 200


def scanned_document_warning(file_path: str, text: str):
    """A plain-English warning when a file yielded almost no readable text, else None.

    CONFIRMED REAL CASE (product owner): a rate table sent as a SCREENSHOT saved to PDF. The
    whole page is one image, so text extraction returns 15 characters - and the app carried on
    as though it had read the document. Everything downstream then works from nothing: detection
    finds no products, or extraction fills a screen with plausible blanks, and the only clue is
    that the answers are wrong in a way that looks like the AI being stupid rather than the AI
    being handed an empty page.

    Cheap to detect, and impossible to diagnose from the symptoms - so it is said out loud."""
    if len((text or "").strip()) >= _MIN_USEFUL_CHARS:
        # CONFIRMED BUG FIX (audit 2026-09-01, MEDIUM/LOW batch 2): this whole-document check
        # passed silently for a real, verified failure mode - a cover page of genuine text
        # (easily clearing 200 characters on its own) followed by pricing pages that are pure
        # screenshots. The document as a whole looked "readable"; every page that actually
        # mattered for pricing was blank. Per-page pages still need their own check even when
        # the concatenated total is fine.
        return _scanned_pages_warning(file_path, text)
    name = os.path.basename(file_path or "this file")
    extension = os.path.splitext(file_path or "")[1].lower()
    if extension == ".pdf":
        return (f"**{name} contains almost no readable text** ({len((text or '').strip())} "
                f"characters). It is almost certainly a scan or a screenshot - a picture of a "
                f"table rather than a table. Nothing can be extracted from it reliably. Please "
                f"upload the original Word or Excel file, or a PDF exported from it, rather than "
                f"a photographed or screenshotted page.")
    return (f"**{name} contains almost no readable text** ({len((text or '').strip())} "
            f"characters). If the content is inside images, the app cannot read it - please "
            f"send the underlying table instead.")


_MIN_USEFUL_CHARS_PER_PAGE = 40


def _scanned_pages_warning(file_path: str, text: str):
    """A plain-English warning naming specific near-blank PAGES, when the document as a whole
    passed scanned_document_warning's total-length check but one or more individual pages did
    not - e.g. a text cover page followed by screenshot pricing pages. Only meaningful for PDFs,
    the only format extract_text_from_pdf marks with "--- Page N ---" per-page breaks; docx/xlsx
    have no comparable per-page concept."""
    extension = os.path.splitext(file_path or "")[1].lower()
    if extension != ".pdf" or "--- Page " not in (text or ""):
        return None
    pages = (text or "").split("--- Page ")[1:]
    blank_pages = []
    for page_block in pages:
        header, _, body = page_block.partition(" ---\n")
        page_num = header.split(" ", 1)[0].strip()
        if len(body.strip()) < _MIN_USEFUL_CHARS_PER_PAGE:
            blank_pages.append(page_num)
    if not blank_pages or len(blank_pages) == len(pages):
        # All pages blank is already caught by the whole-document check above; only a genuine
        # MIX (some fine, some blank) is new information this check adds.
        return None
    name = os.path.basename(file_path or "this file")
    page_word = "page" if len(blank_pages) == 1 else "pages"
    return (f"**{name}: {page_word} {', '.join(blank_pages)} contain almost no readable text**, "
            f"even though the document as a whole does. Those pages are almost certainly a scan "
            f"or a screenshot - if pricing lives there, it will come back blank. Please check "
            f"those pages, or upload the original file rather than a photographed/screenshotted "
            f"one.")


def extract_raw_text(file_path: str) -> str:
    """
    Dispatches to the right extractor based on file extension.
    Raises ValueError with a clear message for unsupported formats
    (e.g. legacy .doc / .xls - these need converting to .docx / .xlsx first).
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext in (".doc", ".xls", ".ppt"):
        raise ValueError(
            f"Legacy '{ext}' format isn't supported directly. "
            f"Please re-save/export the file as '.docx', '.xlsx', or '.pptx' first."
        )
    else:
        raise ValueError(f"Unsupported file type: '{ext}'. Supported: .pdf, .docx, .xlsx, .pptx")


import hashlib


def extract_images_from_pdf(file_path: str, max_images: int = 12, seen_hashes: set = None) -> list:
    """Returns list of (image_bytes, extension) tuples for embedded images in a PDF.
    Skips duplicate images (e.g. a logo repeated on every page) via content hash."""
    import fitz  # PyMuPDF

    if seen_hashes is None:
        seen_hashes = set()
    images = []
    doc = fitz.open(file_path)
    try:
        for page_index in range(len(doc)):
            if len(images) >= max_images:
                break
            page = doc[page_index]
            for img in page.get_images(full=True):
                if len(images) >= max_images:
                    break
                xref = img[0]
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                img_hash = hashlib.sha256(img_bytes).hexdigest()
                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)
                images.append((img_bytes, base_image["ext"]))
    finally:
        doc.close()
    return images


def extract_images_from_docx(file_path: str, max_images: int = 12, seen_hashes: set = None) -> list:
    """Returns list of (image_bytes, extension) tuples for embedded images in a Word doc.
    Skips duplicate images via content hash."""
    import docx

    if seen_hashes is None:
        seen_hashes = set()
    doc = docx.Document(file_path)
    images = []
    for rel in doc.part.rels.values():
        if len(images) >= max_images:
            break
        if "image" in rel.reltype:
            img_bytes = rel.target_part.blob
            img_hash = hashlib.sha256(img_bytes).hexdigest()
            if img_hash in seen_hashes:
                continue
            seen_hashes.add(img_hash)
            ext = rel.target_ref.rsplit(".", 1)[-1] if "." in rel.target_ref else "png"
            images.append((img_bytes, ext))
    return images


def extract_images_from_xlsx(file_path: str, max_images: int = 12, seen_hashes: set = None) -> list:
    """Returns list of (image_bytes, extension) tuples for embedded images in an Excel file.
    Skips duplicate images via content hash."""
    import openpyxl

    if seen_hashes is None:
        seen_hashes = set()
    wb = openpyxl.load_workbook(file_path)
    images = []
    for ws in wb.worksheets:
        if len(images) >= max_images:
            break
        for img in getattr(ws, "_images", []):
            if len(images) >= max_images:
                break
            try:
                data = img._data()
                img_hash = hashlib.sha256(data).hexdigest()
                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)
                images.append((data, "png"))
            except Exception:
                continue
    return images


def extract_images_from_pptx(file_path: str, max_images: int = 12, seen_hashes: set = None) -> list:
    """Returns list of (image_bytes, extension) tuples for images embedded in a PowerPoint deck.
    Skips duplicate images via content hash, and walks into grouped shapes - a supplier deck
    commonly groups a photo together with a caption or logo overlay, and python-pptx only
    iterates top-level shapes by default.

    CONFIRMED REAL GAP, found testing against Chris's actual "1 Day Diving Course.pptx"
    (2026-09-04): two separate problems, both on the SAME real file, each hiding the other.

    1) The deck's real hero photos (5 images, ~80KB-1.1MB each) were NOT plain
       MSO_SHAPE_TYPE.PICTURE shapes - they were "Picture Placeholder" shapes
       (MSO_SHAPE_TYPE.PLACEHOLDER, whose placeholder_format.type is PICTURE), a normal pattern
       in a polished slide template. An earlier version of this function only matched
       MSO_SHAPE_TYPE.PICTURE and so found 12 tiny decorative icons (600-3300 bytes each) while
       silently missing every real photo in the deck. Fixed by trying shape.image on every
       non-group shape rather than gating on shape_type first - shape.image raises (caught
       below) for any shape that isn't actually holding a picture, including an empty, unfilled
       placeholder, so this is a broader net that also matches picture placeholders and any
       other picture-holding shape type without hand-listing each one.

    2) Even with (1) fixed, the SAME deck's decorative icons (there were 12 of them, sitting
       earlier in shape order than the 5 real photos) still filled the default max_images=12 cap
       before the walk ever reached a real photo - a genuine 0-real-images-returned result that
       would have looked identical to "this deck has no images", exactly the silent-failure
       shape this whole pipeline has been bitten by before (see the 2026-08-30/08-31 incident
       docs). Real tour/product photos are essentially always far larger than a decorative
       icon/logo/divider graphic, so instead of capping during the walk (which just keeps
       whatever was found first), every image on the slide(s) is collected, then sorted largest
       first before max_images is applied - the biggest, most likely-to-be-real images win a
       slot over small decorative graphics regardless of where they sit in the shape tree."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx import Presentation

    if seen_hashes is None:
        seen_hashes = set()
    candidates = []  # every deduped image found, unranked - sorted/capped after the full walk

    def _walk(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _walk(shape.shapes)
                continue
            try:
                image = shape.image
            except Exception:
                continue
            try:
                img_bytes = image.blob
                img_hash = hashlib.sha256(img_bytes).hexdigest()
                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)
                candidates.append((img_bytes, image.ext))
            except Exception:
                continue

    prs = Presentation(file_path)
    for slide in prs.slides:
        _walk(slide.shapes)
    candidates.sort(key=lambda pair: len(pair[0]), reverse=True)
    return candidates[:max_images]


def extract_images(file_path: str, max_images: int = 12, seen_hashes: set = None,
                    errors: list = None, label: str = None) -> list:
    """
    Dispatches to the right image extractor based on file extension.
    Returns list of (image_bytes, extension) tuples, or an empty list if
    the format isn't supported for image extraction or none are found.
    seen_hashes: an optional shared set of content hashes to dedupe against
    - pass the SAME set across multiple calls (e.g. multiple uploaded
    documents in one session) to dedupe across files too, not just within one.

    errors: optional list to append a human-readable message to when image extraction genuinely
    FAILED (a corrupted file, a PyMuPDF/python-docx/openpyxl internal error, an unsupported file
    type) - as opposed to the document simply having no embedded images, which stays silent (not
    an error, just nothing to report). CONFIRMED REAL BUG (reported, 2026-08-31: "the App never
    even shows me available images, even though the document...has some images included") - a
    genuine extraction failure used to only reach a print() statement, which never reaches a
    human on Streamlit Cloud (stdout isn't shown anywhere in the deployed app's UI), so it looked
    EXACTLY like "this document has no images" - the same class of bug already fixed once for the
    R2 upload step (see upload_images_r2_with_errors) but left unfixed one step earlier, at
    extraction itself. Left as an opt-in output parameter, not a return-value change, so every
    existing direct caller (all four call sites in app.py, none of which pass it today) keeps
    working unchanged until they're updated to show it.

    label: the name to use in an error message - defaults to file_path's own basename, but
    file_path is usually a throwaway tempfile path (e.g. "/tmp/xyz123.pdf") by the time this
    runs, so callers should pass the ORIGINAL uploaded filename here for a message a human can
    actually recognize.
    """
    if seen_hashes is None:
        seen_hashes = set()
    display_name = label or os.path.basename(file_path)
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            return extract_images_from_pdf(file_path, max_images, seen_hashes)
        elif ext == ".docx":
            return extract_images_from_docx(file_path, max_images, seen_hashes)
        elif ext == ".xlsx":
            return extract_images_from_xlsx(file_path, max_images, seen_hashes)
        elif ext == ".pptx":
            return extract_images_from_pptx(file_path, max_images, seen_hashes)
        # No else/error here for an unsupported extension (.doc, .xls, etc.) - that's a normal,
        # expected case (image extraction is a bonus, not a requirement, for a format that isn't
        # PDF/.docx/.xlsx/.pptx), not a failure worth flagging.
    except Exception as e:
        print(f"⚠️ Image extraction failed for {file_path}: {e}")
        if errors is not None:
            errors.append(f"'{display_name}': couldn't read embedded images from this file - {e}")
    return []
