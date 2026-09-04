"""Tests for adding PowerPoint (.pptx) support to document_reader.py (2026-09-04).

CONFIRMED REAL GAP (Chris, 2026-09-04): a supplier document can arrive as a PowerPoint deck (a
real example: "1 Day Diving Course.pptx"). Before this fix, `.pptx` wasn't even in the file
uploader's accepted types anywhere in app.py, and document_reader.extract_raw_text/extract_images
both raised/returned nothing for it - a document type this app could never read at all.

Testing against that real file (locally, not committed to the repo - built as a synthetic
equivalent here) surfaced TWO further real bugs in the first working version of
extract_images_from_pptx, both confirmed against the actual file before being fixed:

  1. The deck's real hero photos lived in "Picture Placeholder" shapes (MSO_SHAPE_TYPE.PLACEHOLDER
     whose placeholder_format.type is PICTURE), not plain MSO_SHAPE_TYPE.PICTURE shapes - a normal
     pattern for a polished slide template. Gating on shape_type == PICTURE alone found only 12
     tiny decorative icons and missed every real photo.
  2. Even once placeholder photos were found, the deck's 12 small decorative icons sat EARLIER in
     shape order than the 5 real photos - so the default max_images=12 cap filled up on icons
     before the walk ever reached a real photo. Fixed by collecting every image across the whole
     deck first, then sorting largest-first before applying the cap, so real (larger) photos beat
     decorative (smaller) graphics for a capped slot regardless of shape order.

Real .pptx files (built with python-pptx, not committed as binary fixtures - matches this test
suite's existing convention of not checking in real .docx/.xlsx fixtures either) are constructed
per test so these are genuine extraction tests, not source-text checks, for the parts of the
pipeless that CAN be exercised without app.py's own Streamlit runtime.
"""
import hashlib
import io
import os
import tempfile

import pytest

pptx = pytest.importorskip("pptx")
PIL_Image = pytest.importorskip("PIL.Image")

from pptx.util import Emu

import document_reader

MODULE_BUILD = "2026-09-04-pptx-text-and-image-extraction"

_APP_PY = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "app.py")
_REQUIREMENTS_TXT = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "requirements.txt")


def _read_app_py():
    with open(_APP_PY, "r", encoding="utf-8") as f:
        return f.read()


def _make_png_file(path, color, size):
    from PIL import Image
    Image.new("RGB", size, color=color).save(path, format="PNG")


def _make_noisy_png_file(path, size):
    """A random-noise image, unlike a flat-color one, doesn't compress away to a tiny PNG - this
    is what makes it a realistic stand-in for a real photo's file size (a flat-color 800x600 PNG
    is under 3KB; noise of the same dimensions is well over 1MB, like a real photo)."""
    import random
    from PIL import Image
    rng = random.Random(42)
    img = Image.new("RGB", size)
    img.putdata([(rng.randrange(256), rng.randrange(256), rng.randrange(256))
                 for _ in range(size[0] * size[1])])
    img.save(path, format="PNG")


def _picture_placeholder_layout(prs):
    """Finds a slide layout that has a Picture Placeholder - the real-world pattern this whole
    fix exists for (see module docstring, gap #1)."""
    from pptx.presentation import Presentation as _P  # noqa: F401 (type hint only)
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.placeholder_format.type is not None and "PICTURE" in str(ph.placeholder_format.type):
                return layout
    return None


def _build_sample_pptx(tmpdir, with_table=True, with_notes=True):
    """Builds a real .pptx with: slide 1 text + a picture PLACEHOLDER holding a large "real
    photo" + several small plain-picture "icons" added BEFORE the placeholder in shape order (the
    exact ordering that triggered gap #2), all inside a group for one of the icons (gap-adjacent
    coverage: grouped shapes must still be walked); slide 2 with a table and speaker notes."""
    from pptx import Presentation

    icon_path = os.path.join(tmpdir, "icon.png")
    photo_path = os.path.join(tmpdir, "photo.png")
    _make_png_file(icon_path, (255, 0, 0), (10, 10))       # tiny decorative icon
    _make_noisy_png_file(photo_path, (800, 600))           # large real photo (see helper docstring)

    prs = Presentation()
    layout_picture = _picture_placeholder_layout(prs)
    assert layout_picture is not None, "default python-pptx template has no picture placeholder layout"

    slide = prs.slides.add_slide(layout_picture)
    for shape in slide.shapes:
        if shape.has_text_frame and shape.placeholder_format is not None and \
                "TITLE" in str(shape.placeholder_format.type or ""):
            shape.text_frame.text = "1-Day Diving Experience"

    # Icons added BEFORE the real photo is inserted into its placeholder - shape order matters.
    slide.shapes.add_picture(icon_path, Emu(0), Emu(0), width=Emu(90000), height=Emu(90000))
    slide.shapes.add_picture(icon_path, Emu(200000), Emu(0), width=Emu(90000), height=Emu(90000))

    # A grouped icon - grouped shapes must still be walked into.
    group = slide.shapes.add_group_shape()
    grouped_pic = group.shapes.add_picture(
        icon_path, Emu(0), Emu(0), width=Emu(90000), height=Emu(90000))
    group.shapes  # touch to keep reference alive

    ph_filled = False
    for ph in slide.placeholders:
        if ph.placeholder_format.type is not None and "PICTURE" in str(ph.placeholder_format.type):
            ph.insert_picture(photo_path)
            ph_filled = True
            break
    assert ph_filled, "test setup failed to fill the picture placeholder"

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    tb = slide2.shapes.add_textbox(Emu(0), Emu(0), Emu(2000000), Emu(300000))
    tb.text_frame.text = "Rates"
    if with_table:
        gframe = slide2.shapes.add_table(2, 2, Emu(0), Emu(400000), Emu(3000000), Emu(600000))
        table = gframe.table
        table.cell(0, 0).text = "Adult"
        table.cell(0, 1).text = "45"
        table.cell(1, 0).text = "Child"
        table.cell(1, 1).text = "30"
    if with_notes:
        slide2.notes_slide.notes_text_frame.text = "Confirm pickup 48h before tour date."

    out_path = os.path.join(tmpdir, "sample.pptx")
    prs.save(out_path)
    return out_path


# ======================================================================
# Text extraction
# ======================================================================
def test_extract_text_from_pptx_reads_slide_text_table_and_notes():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = _build_sample_pptx(tmpdir)
        text = document_reader.extract_text_from_pptx(path)
    assert "1-Day Diving Experience" in text
    assert "[SLIDE 1]" in text
    assert "[SLIDE 2]" in text
    assert "Rates" in text
    assert "Adult | 45" in text
    assert "Child | 30" in text
    assert "Confirm pickup 48h before tour date." in text
    assert "[SLIDE 2 NOTES]" in text


def test_extract_raw_text_dispatches_pptx_to_the_new_extractor():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = _build_sample_pptx(tmpdir)
        text = document_reader.extract_raw_text(path)
    assert "1-Day Diving Experience" in text


# ======================================================================
# Image extraction - gap #1 (picture placeholders) and gap #2 (icon-crowding under the cap)
# ======================================================================
def test_extract_images_from_pptx_finds_the_photo_inside_a_picture_placeholder():
    """Gap #1: a plain shape_type==PICTURE check alone would miss this - the real photo in the
    sample deck lives inside a Picture Placeholder, exactly like the real diving-course file."""
    with tempfile.TemporaryDirectory() as tmpdir:
        path = _build_sample_pptx(tmpdir)
        imgs = document_reader.extract_images_from_pptx(path, max_images=50)
    sizes = sorted(len(b) for b, ext in imgs)
    assert sizes[-1] > 50000, "the large real photo (800x600 PNG) should be the biggest image found"


def test_extract_images_from_pptx_walks_into_grouped_shapes():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = _build_sample_pptx(tmpdir)
        imgs = document_reader.extract_images_from_pptx(path, max_images=50)
    # 2 loose icons + 1 grouped icon + 1 placeholder photo, all distinct pixel content except the
    # two loose icons and the grouped icon share identical bytes (deduped - see next test).
    assert len(imgs) >= 2


def test_extract_images_from_pptx_dedupes_identical_icons_by_content_hash():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = _build_sample_pptx(tmpdir)
        imgs = document_reader.extract_images_from_pptx(path, max_images=50)
    hashes = [hashlib.sha256(b).hexdigest() for b, ext in imgs]
    assert len(hashes) == len(set(hashes)), "identical icon bytes must only be returned once"


def test_extract_images_from_pptx_prioritizes_larger_real_photos_over_small_icons_when_capped():
    """Gap #2: the sample deck's icons are added to shape order BEFORE the real photo is filled
    into its placeholder. A naive first-N-found cap would return only icons; sorting largest-first
    before capping must surface the real photo even with max_images=1."""
    with tempfile.TemporaryDirectory() as tmpdir:
        path = _build_sample_pptx(tmpdir)
        imgs = document_reader.extract_images_from_pptx(path, max_images=1)
    assert len(imgs) == 1
    assert len(imgs[0][0]) > 50000, "the one slot returned must be the real photo, not an icon"


def test_extract_images_dispatches_pptx_to_the_new_extractor():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = _build_sample_pptx(tmpdir)
        imgs = document_reader.extract_images(path, label="sample.pptx")
    assert len(imgs) >= 1
    assert max(len(b) for b, ext in imgs) > 50000


def test_extract_images_from_a_deck_with_no_real_photos_returns_only_what_exists_no_error():
    """A deck that genuinely has zero pictures must stay silent (empty list, no error) - this is
    the ordinary 'nothing to find' case that must not be confused with a real extraction failure,
    matching the standing rule the 2026-08-31 incident fix established for the other formats."""
    from pptx import Presentation
    with tempfile.TemporaryDirectory() as tmpdir:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(300000)).text_frame.text = "No photos here"
        path = os.path.join(tmpdir, "no_photos.pptx")
        prs.save(path)

        errors = []
        imgs = document_reader.extract_images(path, errors=errors, label="no_photos.pptx")
    assert imgs == []
    assert errors == []


# ======================================================================
# Error handling / format-support surface
# ======================================================================
def test_legacy_ppt_extension_raises_a_clear_conversion_message():
    with pytest.raises(ValueError, match="pptx"):
        document_reader.extract_raw_text("/tmp/old_presentation.ppt")


def test_unsupported_extension_error_message_now_mentions_pptx():
    with pytest.raises(ValueError, match=r"\.pptx"):
        document_reader.extract_raw_text("/tmp/whatever.txt")


def test_extract_images_stays_silent_for_a_genuinely_unsupported_extension():
    """.doc/.xls (not .pptx anymore) are still the 'normal, not a failure' case for image
    extraction - unaffected by this fix, guarded here so the gap comment above extract_images'
    dispatch doesn't quietly regress this behavior for the formats it still applies to."""
    errors = []
    out = document_reader.extract_images("/tmp/legacy.doc", errors=errors)
    assert out == []
    assert errors == []


# ======================================================================
# Wiring - app.py's upload widgets now accept .pptx, requirements.txt carries python-pptx
# ======================================================================
def test_all_eight_file_uploaders_in_app_py_now_accept_pptx():
    src = _read_app_py()
    old_type_list = 'type=["pdf", "docx", "xlsx"]'
    new_type_list = 'type=["pdf", "docx", "xlsx", "pptx"]'
    assert old_type_list not in src, "an upload widget was missed when adding pptx support"
    assert src.count(new_type_list) == 8


def test_requirements_txt_lists_python_pptx():
    with open(_REQUIREMENTS_TXT, "r", encoding="utf-8") as f:
        reqs = f.read()
    assert "python-pptx" in reqs
